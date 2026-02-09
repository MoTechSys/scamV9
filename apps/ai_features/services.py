"""
خدمات الذكاء الاصطناعي - Google Gemini (Enterprise Edition v2)
S-ACM - Smart Academic Content Management System

=== Phase 2: TANK AI Engine (Dynamic Governance) ===
1. HydraKeyManager: DB-based Round-Robin with cooldown, RPM enforcement
2. SmartChunker: Uses AIConfiguration.chunk_size from DB
3. GeminiService: Uses AIConfiguration for model/tokens/temperature
4. All settings are Admin-editable, ZERO .env dependency for AI config

== Legacy Compatibility ==
- All public interfaces preserved (GeminiService, QuestionMatrixConfig, etc.)
- Falls back to .env keys if no DB keys configured
"""

from __future__ import annotations

import json
import hashlib
import logging
import os
import re
import time
import threading
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from enum import Enum
from functools import wraps
from pathlib import Path
from typing import Optional, List, Dict, Any, Callable, TypeVar
from datetime import datetime

from django.conf import settings
from django.core.cache import cache
from django.utils import timezone

# ========== Logging ==========
logger = logging.getLogger('ai_features')


# ========== Constants (Fallbacks - DB config takes priority) ==========
FALLBACK_MODEL = getattr(settings, 'AI_MODEL_NAME', None) or os.getenv('AI_MODEL_NAME', 'gemini-2.5-flash')
FALLBACK_CHUNK_SIZE = 30000
FALLBACK_CHUNK_OVERLAP = 500
FALLBACK_MAX_OUTPUT_TOKENS = 2000
FALLBACK_TEMPERATURE = 0.3
CACHE_TIMEOUT = 3600
MAX_RETRIES = 3
AI_OUTPUT_DIR = 'ai_generated'

# Legacy compatibility aliases
GEMINI_MODEL = FALLBACK_MODEL
MAX_INPUT_LENGTH = FALLBACK_CHUNK_SIZE
CHUNK_SIZE = 8000
CHUNK_OVERLAP = 500


# ========== Custom Exceptions ==========

class GeminiError(Exception):
    """Base exception for Gemini-related errors."""
    pass


class GeminiConfigurationError(GeminiError):
    """Raised when Gemini is not properly configured."""
    pass


class GeminiAPIError(GeminiError):
    """Raised when Gemini API returns an error."""
    def __init__(self, message: str, status_code: Optional[int] = None):
        super().__init__(message)
        self.status_code = status_code


class GeminiRateLimitError(GeminiAPIError):
    """Raised when rate limit is exceeded."""
    pass


class GeminiServiceDisabledError(GeminiError):
    """Raised when AI service is disabled by admin."""
    pass


class TextExtractionError(GeminiError):
    """Raised when text extraction from file fails."""
    pass


# ========== Enums ==========

class QuestionType(Enum):
    MCQ = "mcq"
    TRUE_FALSE = "true_false"
    SHORT_ANSWER = "short_answer"
    MIXED = "mixed"


class ContentType(Enum):
    PDF = "pdf"
    DOCX = "docx"
    PPTX = "pptx"
    TEXT = "text"
    UNKNOWN = "unknown"


# ========== Data Classes ==========

@dataclass
class Question:
    """نموذج سؤال واحد."""
    type: str
    question: str
    answer: str
    options: Optional[List[str]] = None
    explanation: Optional[str] = None
    score: float = 1.0

    def to_dict(self) -> Dict[str, Any]:
        result = {
            'type': self.type,
            'question': self.question,
            'answer': self.answer,
            'score': self.score,
        }
        if self.options:
            result['options'] = self.options
        if self.explanation:
            result['explanation'] = self.explanation
        return result


@dataclass
class QuestionMatrixConfig:
    """تكوين مصفوفة الأسئلة."""
    mcq_count: int = 0
    mcq_score: float = 2.0
    true_false_count: int = 0
    true_false_score: float = 1.0
    short_answer_count: int = 0
    short_answer_score: float = 3.0

    @property
    def total_questions(self) -> int:
        return self.mcq_count + self.true_false_count + self.short_answer_count

    @property
    def total_score(self) -> float:
        return (
            self.mcq_count * self.mcq_score +
            self.true_false_count * self.true_false_score +
            self.short_answer_count * self.short_answer_score
        )

    def to_dict(self) -> Dict[str, Any]:
        return {
            'mcq_count': self.mcq_count, 'mcq_score': self.mcq_score,
            'true_false_count': self.true_false_count, 'true_false_score': self.true_false_score,
            'short_answer_count': self.short_answer_count, 'short_answer_score': self.short_answer_score,
            'total_questions': self.total_questions, 'total_score': self.total_score,
        }


@dataclass
class AIResponse:
    """نموذج استجابة AI."""
    success: bool
    data: Any = None
    error: Optional[str] = None
    cached: bool = False
    md_file_path: Optional[str] = None


# ========================================================================
# Phase 2: HYDRA KEY MANAGER (DB-based)
# ========================================================================

class HydraKeyManager:
    """
    The "Hydra" Key Manager - DB-Powered Round-Robin with Cooldown.

    Features:
    - Fetches active keys from DB (APIKey model)
    - Round-Robin rotation across available keys
    - Automatic cooldown on 429 (Rate Limit) errors
    - RPM (Requests Per Minute) enforcement per key
    - Falls back to .env keys if no DB keys exist
    - Thread-safe singleton

    Usage:
        manager = HydraKeyManager()
        key_obj, raw_key = manager.get_next_key()
        # ... use raw_key for API call ...
        key_obj.mark_success(latency_ms=150)
        # or on error:
        key_obj.mark_error("429 Too Many Requests", is_rate_limit=True)
    """
    _instance = None
    _lock = threading.Lock()

    def __new__(cls):
        if cls._instance is None:
            with cls._lock:
                if cls._instance is None:
                    cls._instance = super().__new__(cls)
                    cls._instance._initialized = False
        return cls._instance

    def __init__(self):
        if self._initialized:
            return
        self._current_index = 0
        self._key_lock = threading.Lock()
        self._fallback_keys: List[str] = []
        self._fallback_index = 0
        self._load_fallback_keys()
        self._initialized = True

    def _load_fallback_keys(self):
        """Load fallback keys from .env (legacy support)."""
        keys = []
        primary_key = getattr(settings, 'GEMINI_API_KEY', '') or os.getenv('GEMINI_API_KEY', '')
        if primary_key and primary_key != 'your_gemini_api_key_here':
            keys.append(primary_key)
        for i in range(1, 11):
            key = os.getenv(f'GEMINI_API_KEY_{i}', '')
            if key and key not in keys:
                keys.append(key)
        self._fallback_keys = keys
        if keys:
            logger.info(f"HydraKeyManager: Loaded {len(keys)} fallback .env key(s)")

    def _get_db_keys(self):
        """Fetch active & available keys from DB."""
        try:
            from apps.ai_features.models import APIKey
            now = timezone.now()
            return list(
                APIKey.objects.filter(
                    provider='gemini',
                    is_active=True,
                ).exclude(
                    status__in=['disabled', 'error']
                ).order_by('priority', 'pk')
            )
        except Exception as e:
            logger.warning(f"HydraKeyManager: Cannot fetch DB keys: {e}")
            return []

    def get_next_key(self):
        """
        Get the next available API key using Round-Robin.

        Returns:
            tuple: (APIKey_instance_or_None, raw_key_string)
            If using DB key: (APIKey, decrypted_key)
            If using fallback: (None, env_key)

        Raises:
            GeminiConfigurationError: If no keys are available.
        """
        with self._key_lock:
            # Try DB keys first
            db_keys = self._get_db_keys()
            available_keys = [k for k in db_keys if k.is_available() and k.check_rpm_limit()]

            if available_keys:
                key_obj = available_keys[self._current_index % len(available_keys)]
                self._current_index += 1
                raw_key = key_obj.get_key()
                if raw_key:
                    return key_obj, raw_key
                logger.warning(f"HydraKeyManager: Key {key_obj.label} decryption failed")

            # Fallback to .env keys
            if self._fallback_keys:
                key = self._fallback_keys[self._fallback_index % len(self._fallback_keys)]
                self._fallback_index += 1
                logger.info("HydraKeyManager: Using fallback .env key")
                return None, key

            raise GeminiConfigurationError(
                "لا توجد مفاتيح API متاحة. أضف مفاتيح من لوحة الإدارة أو في ملف .env"
            )

    def rotate_after_error(self, failed_key_obj=None, error_msg: str = '', is_rate_limit: bool = False):
        """Rotate to next key after an error."""
        if failed_key_obj:
            failed_key_obj.mark_error(error_msg, is_rate_limit=is_rate_limit)
        with self._key_lock:
            self._current_index += 1
            if not failed_key_obj and self._fallback_keys:
                self._fallback_index += 1

    @property
    def total_keys(self) -> int:
        db_count = len(self._get_db_keys())
        return db_count + len(self._fallback_keys)

    @property
    def has_keys(self) -> bool:
        return self.total_keys > 0

    def get_health_status(self) -> List[Dict[str, Any]]:
        """Get health status of all keys (for Admin Dashboard)."""
        result = []
        db_keys = self._get_db_keys()
        for key in db_keys:
            result.append({
                'id': key.pk,
                'label': key.label,
                'hint': key.key_hint,
                'status': key.status,
                'is_available': key.is_available(),
                'error_count': key.error_count,
                'total_requests': key.total_requests,
                'last_latency_ms': key.last_latency_ms,
                'last_success_at': key.last_success_at,
                'last_error': key.last_error,
                'rpm_limit': key.rpm_limit,
                'tokens_used_today': key.tokens_used_today,
                'cooldown_until': key.cooldown_until,
            })
        return result


# Legacy compatibility: Alias
APIKeyManager = HydraKeyManager


# ========================================================================
# Smart Text Chunking (DB-Configured)
# ========================================================================

class SmartChunker:
    """
    تقسيم ذكي للنصوص الكبيرة.
    Reads chunk_size and overlap from AIConfiguration (DB).
    """

    def __init__(self, chunk_size: int = None, overlap: int = None):
        if chunk_size is None or overlap is None:
            try:
                from apps.ai_features.models import AIConfiguration
                config = AIConfiguration.get_config()
                self.chunk_size = chunk_size or config.chunk_size
                self.overlap = overlap or config.chunk_overlap
            except Exception:
                self.chunk_size = chunk_size or FALLBACK_CHUNK_SIZE
                self.overlap = overlap or FALLBACK_CHUNK_OVERLAP
        else:
            self.chunk_size = chunk_size
            self.overlap = overlap

    def chunk_text(self, text: str) -> List[str]:
        """تقسيم النص إلى أجزاء ذكية."""
        if not text:
            return []
        if len(text) <= self.chunk_size:
            return [text]

        chunks = []
        paragraphs = text.split('\n\n')
        current_chunk = ""

        for para in paragraphs:
            para = para.strip()
            if not para:
                continue

            if len(para) > self.chunk_size:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                    current_chunk = ""
                sentence_chunks = self._split_by_sentences(para)
                chunks.extend(sentence_chunks)
                continue

            if len(current_chunk) + len(para) + 2 > self.chunk_size:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                    overlap_text = current_chunk[-self.overlap:] if len(current_chunk) > self.overlap else ""
                    current_chunk = overlap_text + "\n\n" + para
                else:
                    current_chunk = para
            else:
                current_chunk += ("\n\n" if current_chunk else "") + para

        if current_chunk.strip():
            chunks.append(current_chunk.strip())

        logger.info(f"SmartChunker: Split {len(text)} chars into {len(chunks)} chunks (size={self.chunk_size})")
        return chunks

    def _split_by_sentences(self, text: str) -> List[str]:
        """تقسيم النص بالجمل."""
        separators = ['. ', '.\n', '。', '؟ ', '? ', '! ', '！ ', '.\t']
        sentences = [text]
        for sep in separators:
            new_sentences = []
            for s in sentences:
                parts = s.split(sep)
                for i, part in enumerate(parts):
                    if i < len(parts) - 1:
                        new_sentences.append(part + sep.strip())
                    else:
                        if part.strip():
                            new_sentences.append(part)
            sentences = new_sentences

        chunks = []
        current = ""
        for sentence in sentences:
            if len(current) + len(sentence) + 1 > self.chunk_size:
                if current:
                    chunks.append(current.strip())
                current = sentence
            else:
                current += " " + sentence if current else sentence

        if current.strip():
            chunks.append(current.strip())

        return chunks


# ========== File-Based AI Storage ==========

class AIFileStorage:
    """مخزن ملفات AI - يحفظ المخرجات كملفات .md في media/ai_generated/."""

    def __init__(self):
        self.base_dir = Path(settings.MEDIA_ROOT) / AI_OUTPUT_DIR
        self.base_dir.mkdir(parents=True, exist_ok=True)

    def save_summary(self, file_id: int, content: str, metadata: Optional[Dict] = None) -> str:
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"summary_{file_id}_{timestamp}.md"
        return self._save_file(filename, content, metadata, "summary")

    def save_questions(self, file_id: int, questions_data: List[Dict], metadata: Optional[Dict] = None) -> str:
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"questions_{file_id}_{timestamp}.md"
        md_content = self._questions_to_markdown(questions_data, metadata)
        return self._save_file(filename, md_content, None, "questions")

    def save_chat_answer(self, file_id: int, user_id: int, question: str, answer: str) -> str:
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"chat_{file_id}_{user_id}_{timestamp}.md"
        content = f"# سؤال\n\n{question}\n\n# إجابة\n\n{answer}\n"
        return self._save_file(filename, content, None, "chat")

    def read_file(self, relative_path: str) -> Optional[str]:
        full_path = Path(settings.MEDIA_ROOT) / relative_path
        try:
            if full_path.exists():
                return full_path.read_text(encoding='utf-8')
        except Exception as e:
            logger.error(f"AIFileStorage: Failed to read {full_path}: {e}")
        return None

    def delete_file(self, relative_path: str) -> bool:
        full_path = Path(settings.MEDIA_ROOT) / relative_path
        try:
            if full_path.exists():
                full_path.unlink()
                logger.info(f"AIFileStorage: Deleted {relative_path}")
                return True
        except Exception as e:
            logger.error(f"AIFileStorage: Failed to delete {relative_path}: {e}")
        return False

    def _save_file(self, filename: str, content: str, metadata: Optional[Dict], category: str) -> str:
        category_dir = self.base_dir / category
        category_dir.mkdir(parents=True, exist_ok=True)

        filepath = category_dir / filename
        header = ""
        if metadata:
            header = "---\n"
            for k, v in metadata.items():
                header += f"{k}: {v}\n"
            header += "---\n\n"

        filepath.write_text(header + content, encoding='utf-8')
        relative_path = str(Path(AI_OUTPUT_DIR) / category / filename)
        logger.info(f"AIFileStorage: Saved {relative_path} ({len(content)} chars)")
        return relative_path

    def _questions_to_markdown(self, questions: List[Dict], metadata: Optional[Dict] = None) -> str:
        lines = ["# بنك الأسئلة المُولَّدة بالذكاء الاصطناعي\n"]

        if metadata:
            lines.append(f"**المصدر:** {metadata.get('source_file', 'غير محدد')}")
            lines.append(f"**التاريخ:** {metadata.get('date', datetime.now().strftime('%Y-%m-%d'))}")
            lines.append(f"**إجمالي الدرجات:** {metadata.get('total_score', '-')}")
            lines.append("")

        type_labels = {
            'mcq': 'اختيار من متعدد',
            'true_false': 'صح وخطأ',
            'short_answer': 'إجابة قصيرة'
        }

        for i, q in enumerate(questions, 1):
            q_type = q.get('type', 'short_answer')
            score = q.get('score', 1.0)
            label = type_labels.get(q_type, q_type)

            lines.append(f"## السؤال {i} ({label}) - [{score} درجة]")
            lines.append(f"\n{q.get('question', '')}\n")

            options = q.get('options')
            if options and isinstance(options, list):
                for j, opt in enumerate(options):
                    letter = chr(ord('أ') + j) if j < 4 else chr(ord('a') + j)
                    lines.append(f"- {letter}) {opt}")
                lines.append("")

            lines.append(f"**الإجابة:** {q.get('answer', '')}")

            explanation = q.get('explanation')
            if explanation:
                lines.append(f"\n**الشرح:** {explanation}")
            lines.append("\n---\n")

        return "\n".join(lines)


# ========== Decorators ==========

T = TypeVar('T')


def cache_result(timeout: int = CACHE_TIMEOUT):
    """Decorator لتخزين نتائج AI في الكاش."""
    def decorator(func: Callable[..., T]) -> Callable[..., T]:
        @wraps(func)
        def wrapper(self, text: str, *args, **kwargs) -> T:
            cache_key = _generate_cache_key(func.__name__, text, args, kwargs)
            cached_result = cache.get(cache_key)
            if cached_result is not None:
                logger.debug(f"Cache hit for {func.__name__}")
                return cached_result
            result = func(self, text, *args, **kwargs)
            if result is not None:
                cache.set(cache_key, result, timeout)
            return result
        return wrapper
    return decorator


def _generate_cache_key(func_name: str, text: str, args: tuple, kwargs: dict) -> str:
    content = f"{func_name}:{text[:200]}:{str(args)}:{str(sorted(kwargs.items()))}"
    return f"ai:{hashlib.md5(content.encode()).hexdigest()}"


# ========== Text Extractors ==========

class TextExtractor(ABC):
    @abstractmethod
    def extract(self, file_path: Path) -> str:
        pass

    @abstractmethod
    def supports(self, file_path: Path) -> bool:
        pass


class PDFExtractor(TextExtractor):
    def supports(self, file_path: Path) -> bool:
        return file_path.suffix.lower() == '.pdf'

    def extract(self, file_path: Path) -> str:
        try:
            import pdfplumber
        except ImportError:
            raise TextExtractionError("pdfplumber not installed. Run: pip install pdfplumber")
        try:
            text_parts = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
            return "\n".join(text_parts)
        except Exception as e:
            raise TextExtractionError(f"Failed to extract text from PDF: {e}")


class DocxExtractor(TextExtractor):
    def supports(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in ['.docx', '.doc']

    def extract(self, file_path: Path) -> str:
        try:
            from docx import Document
        except ImportError:
            raise TextExtractionError("python-docx not installed. Run: pip install python-docx")
        try:
            doc = Document(file_path)
            return "\n".join(para.text for para in doc.paragraphs if para.text)
        except Exception as e:
            raise TextExtractionError(f"Failed to extract text from DOCX: {e}")


class PptxExtractor(TextExtractor):
    def supports(self, file_path: Path) -> bool:
        return file_path.suffix.lower() == '.pptx'

    def extract(self, file_path: Path) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            raise TextExtractionError("python-pptx not installed. Run: pip install python-pptx")
        try:
            prs = Presentation(file_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)
            return "\n".join(text_parts)
        except Exception as e:
            raise TextExtractionError(f"Failed to extract text from PPTX: {e}")


class PlainTextExtractor(TextExtractor):
    SUPPORTED_EXTENSIONS = {'.txt', '.md', '.rst', '.csv'}

    def supports(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in self.SUPPORTED_EXTENSIONS

    def extract(self, file_path: Path) -> str:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            with open(file_path, 'r', encoding='cp1256') as f:
                return f.read()
        except Exception as e:
            raise TextExtractionError(f"Failed to read text file: {e}")


class TextExtractorFactory:
    _extractors: List[TextExtractor] = [
        PDFExtractor(),
        DocxExtractor(),
        PptxExtractor(),
        PlainTextExtractor(),
    ]

    @classmethod
    def get_extractor(cls, file_path: Path) -> Optional[TextExtractor]:
        for extractor in cls._extractors:
            if extractor.supports(file_path):
                return extractor
        return None

    @classmethod
    def extract_text(cls, file_path: Path) -> str:
        extractor = cls.get_extractor(file_path)
        if extractor is None:
            raise TextExtractionError(f"Unsupported file type: {file_path.suffix}")
        return extractor.extract(file_path)


# ========================================================================
# Gemini Service (Enterprise v2 - DB Governed)
# ========================================================================

def _get_ai_config():
    """Helper to get AI configuration from DB with fallback."""
    try:
        from apps.ai_features.models import AIConfiguration
        return AIConfiguration.get_config()
    except Exception:
        return None


class GeminiService:
    """
    خدمة Google Gemini للذكاء الاصطناعي - Enterprise v2.

    === Dynamic Governance ===
    - Model, tokens, temperature: Read from AIConfiguration (DB)
    - API Keys: Managed by HydraKeyManager (DB + .env fallback)
    - Chunk size: Read from AIConfiguration (DB)
    - Rate limits: Per-user (DB) + Per-key RPM (DB)
    - Service toggle: Can be disabled from Admin panel

    === Admin Editable (No Code Touch) ===
    - Change model: Admin -> AI Configuration -> active_model
    - Add keys: Admin -> API Keys -> Add
    - Change RPM: Admin -> API Keys -> rpm_limit
    - Disable service: Admin -> AI Configuration -> is_service_enabled
    """

    def __init__(self, model: str = None):
        config = _get_ai_config()
        self._model_name = model or (config.active_model if config else FALLBACK_MODEL)
        self._key_manager = HydraKeyManager()
        self._chunker = SmartChunker()
        self._storage = AIFileStorage()
        self._client = None
        self._current_key_obj = None
        self._initialize_client()

    def _check_service_enabled(self):
        """Check if AI service is enabled by admin."""
        config = _get_ai_config()
        if config and not config.is_service_enabled:
            msg = config.maintenance_message or 'خدمة الذكاء الاصطناعي متوقفة مؤقتاً.'
            raise GeminiServiceDisabledError(msg)

    def _initialize_client(self) -> None:
        """تهيئة عميل Gemini باستخدام المفتاح الحالي."""
        if not self._key_manager.has_keys:
            logger.warning("GeminiService: No API keys available. Service will be limited.")
            return
        try:
            from google import genai
            key_obj, raw_key = self._key_manager.get_next_key()
            self._current_key_obj = key_obj
            self._client = genai.Client(api_key=raw_key)
            logger.info(f"GeminiService initialized with model: {self._model_name}")
        except ImportError:
            raise GeminiConfigurationError("google-genai not installed. Run: pip install google-genai")
        except GeminiConfigurationError:
            logger.warning("GeminiService: No keys available for initialization.")
        except Exception as e:
            raise GeminiConfigurationError(f"Failed to initialize Gemini client: {e}")

    def _reinitialize_with_next_key(self, error_msg: str = '', is_rate_limit: bool = False) -> None:
        """إعادة تهيئة العميل بالمفتاح التالي."""
        self._key_manager.rotate_after_error(self._current_key_obj, error_msg, is_rate_limit)
        try:
            from google import genai
            key_obj, raw_key = self._key_manager.get_next_key()
            self._current_key_obj = key_obj
            self._client = genai.Client(api_key=raw_key)
            logger.info("GeminiService: Reinitialized with next API key")
        except Exception as e:
            logger.error(f"GeminiService: Failed to reinitialize: {e}")

    @property
    def is_available(self) -> bool:
        return self._client is not None

    @property
    def storage(self) -> AIFileStorage:
        return self._storage

    def _generate_content(self, prompt: str, max_tokens: int = None) -> str:
        """توليد محتوى مع Hydra Round-Robin retry."""
        self._check_service_enabled()

        if not self.is_available:
            raise GeminiConfigurationError("Gemini client not initialized")

        # Get config from DB
        config = _get_ai_config()
        if max_tokens is None:
            max_tokens = config.max_output_tokens if config else FALLBACK_MAX_OUTPUT_TOKENS
        temperature = config.temperature if config else FALLBACK_TEMPERATURE

        last_exception = None
        max_attempts = min(self._key_manager.total_keys, MAX_RETRIES) if self._key_manager.total_keys > 0 else MAX_RETRIES

        for attempt in range(max(max_attempts, 1)):
            try:
                from google import genai
                from google.genai import types

                start_ms = int(time.time() * 1000)

                response = self._client.models.generate_content(
                    model=self._model_name,
                    contents=prompt,
                    config=types.GenerateContentConfig(
                        max_output_tokens=max_tokens,
                        temperature=temperature,
                    )
                )

                latency_ms = int(time.time() * 1000) - start_ms

                if response.text:
                    # Mark success on the current key
                    if self._current_key_obj:
                        self._current_key_obj.mark_success(latency_ms)
                    return response.text.strip()
                else:
                    raise GeminiAPIError("Empty response from Gemini")

            except GeminiServiceDisabledError:
                raise
            except Exception as e:
                last_exception = e
                error_str = str(e).lower()

                is_rate_limit = any(kw in error_str for kw in [
                    "rate", "quota", "429", "resource_exhausted"
                ])

                if is_rate_limit:
                    logger.warning(f"Rate limit on attempt {attempt + 1}, rotating key...")
                    self._reinitialize_with_next_key(str(e), is_rate_limit=True)

                    retry_match = re.search(r'retry in (\d+)', error_str)
                    if retry_match:
                        wait_time = min(int(retry_match.group(1)), 60)
                    else:
                        wait_time = min(5.0 * (2 ** attempt), 30)

                    time.sleep(wait_time)
                    continue
                elif "invalid" in error_str and "key" in error_str:
                    logger.error(f"Invalid API key, rotating...")
                    self._reinitialize_with_next_key(str(e), is_rate_limit=False)
                    continue
                else:
                    if attempt < max_attempts - 1:
                        time.sleep(1.0 * (2 ** attempt))
                        self._reinitialize_with_next_key(str(e), is_rate_limit=False)
                        continue
                    raise GeminiAPIError(f"Gemini API error: {e}")

        if last_exception and any(kw in str(last_exception).lower() for kw in ["quota", "429", "rate"]):
            raise GeminiRateLimitError(
                "⏳ تم تجاوز الحد المسموح لـ API. يرجى الانتظار دقيقة ثم المحاولة مرة أخرى، "
                "أو تواصل مع المسؤول لإضافة مفاتيح API إضافية."
            )
        raise last_exception or GeminiAPIError("All API keys exhausted")

    # ========== Public Methods ==========

    def extract_text_from_file(self, file_obj) -> Optional[str]:
        """استخراج النص من كائن LectureFile."""
        if not file_obj.local_file:
            logger.warning(f"File {file_obj.id} has no local file")
            return None
        try:
            file_path = Path(file_obj.local_file.path)
            text = TextExtractorFactory.extract_text(file_path)
            logger.info(f"Extracted {len(text)} characters from {file_path.name}")
            return text
        except TextExtractionError as e:
            logger.error(f"Text extraction failed for file {file_obj.id}: {e}")
            return None

    @cache_result(timeout=CACHE_TIMEOUT)
    def generate_summary(self, text: str, max_length: int = 500, user_notes: str = "") -> str:
        """توليد تلخيص للنص مع دعم Smart Chunking."""
        chunks = self._chunker.chunk_text(text)

        if len(chunks) <= 1:
            return self._generate_single_summary(text, max_length, user_notes)

        chunk_summaries = []
        for i, chunk in enumerate(chunks):
            try:
                summary = self._generate_single_summary(
                    chunk, max_length=200,
                    user_notes=f"هذا الجزء {i+1} من {len(chunks)}. {user_notes}"
                )
                chunk_summaries.append(summary)
            except GeminiError:
                continue

        if not chunk_summaries:
            return self._fallback_summary(text, max_length)

        combined = "\n\n".join(chunk_summaries)
        return self._generate_single_summary(
            combined, max_length,
            user_notes=f"هذا ملخص مُجمّع من {len(chunks)} أجزاء. أعد صياغته كملخص متماسك واحد. {user_notes}"
        )

    def _generate_single_summary(self, text: str, max_length: int = 500, user_notes: str = "") -> str:
        notes_section = f"\nملاحظات إضافية من المستخدم: {user_notes}" if user_notes else ""
        config = _get_ai_config()
        input_limit = config.chunk_size if config else FALLBACK_CHUNK_SIZE

        prompt = f"""أنت مساعد أكاديمي متخصص في تلخيص المحتوى التعليمي باللغة العربية.

قم بتلخيص النص التالي بشكل مختصر ومفيد بصيغة Markdown. ركز على:
- النقاط الرئيسية والمفاهيم الأساسية
- المعلومات الأكثر أهمية
- الحفاظ على الدقة العلمية
- استخدام عناوين وقوائم لتنظيم المحتوى
{notes_section}

النص:
{text[:input_limit]}

التلخيص (بحد أقصى {max_length} كلمة، بصيغة Markdown):"""

        try:
            return self._generate_content(prompt, max_tokens=max_length * 3)
        except GeminiError as e:
            logger.error(f"Summary generation failed: {e}")
            return self._fallback_summary(text, max_length)

    def _fallback_summary(self, text: str, max_length: int) -> str:
        sentences = text.replace('\n', ' ').split('.')
        summary = ""
        for sentence in sentences:
            sentence = sentence.strip()
            if sentence and len(summary) + len(sentence) < max_length * 4:
                summary += sentence + ". "
            elif len(summary) > 100:
                break
        return summary.strip() or text[:max_length * 4] + "..."

    def generate_questions_matrix(
        self, text: str, matrix: QuestionMatrixConfig, user_notes: str = ""
    ) -> List[Dict[str, Any]]:
        """توليد أسئلة حسب المصفوفة المحددة."""
        if matrix.total_questions == 0:
            return []

        chunks = self._chunker.chunk_text(text)
        config = _get_ai_config()
        input_limit = config.chunk_size if config else FALLBACK_CHUNK_SIZE

        source_text = chunks[0] if chunks else text[:input_limit]
        if len(chunks) > 1:
            source_text = "\n\n---\n\n".join(chunks[:3])

        notes_section = f"\nملاحظات: {user_notes}" if user_notes else ""

        parts = []
        if matrix.mcq_count > 0:
            parts.append(f"- {matrix.mcq_count} سؤال اختيار من متعدد (mcq) - كل سؤال {matrix.mcq_score} درجة")
        if matrix.true_false_count > 0:
            parts.append(f"- {matrix.true_false_count} سؤال صح وخطأ (true_false) - كل سؤال {matrix.true_false_score} درجة")
        if matrix.short_answer_count > 0:
            parts.append(f"- {matrix.short_answer_count} سؤال إجابة قصيرة (short_answer) - كل سؤال {matrix.short_answer_score} درجة")

        matrix_text = "\n".join(parts)

        prompt = f"""أنت مدرس جامعي متخصص في إنشاء أسئلة اختبارية باللغة العربية.

أنشئ الأسئلة التالية بالضبط من النص المقدم:
{matrix_text}

إجمالي: {matrix.total_questions} سؤال | الدرجة الكلية: {matrix.total_score}
{notes_section}

أرجع الإجابة بصيغة JSON فقط بدون أي نص إضافي، كمصفوفة:
[
    {{
        "type": "mcq" أو "true_false" أو "short_answer",
        "question": "نص السؤال",
        "options": ["خيار1", "خيار2", "خيار3", "خيار4"],
        "answer": "الإجابة الصحيحة",
        "explanation": "شرح مختصر",
        "score": الدرجة_كرقم
    }}
]

ملاحظات:
- للأسئلة true_false: الخيارات ["صح", "خطأ"] فقط
- للأسئلة short_answer: لا تضع options (اجعلها null)
- تأكد أن الأسئلة متنوعة وتغطي أجزاء مختلفة من النص

النص:
{source_text[:input_limit]}

الأسئلة (JSON فقط):"""

        try:
            result = self._generate_content(prompt, max_tokens=4000)
            return self._parse_questions_json(result)
        except GeminiError as e:
            logger.error(f"Question matrix generation failed: {e}")
            return self._fallback_questions(matrix.total_questions)

    @cache_result(timeout=CACHE_TIMEOUT)
    def generate_questions(
        self, text: str, question_type: QuestionType = QuestionType.MIXED,
        num_questions: int = 5, user_notes: str = ""
    ) -> List[Dict[str, Any]]:
        """توليد أسئلة (واجهة متوافقة مع الإصدار القديم)."""
        matrix = QuestionMatrixConfig()
        if question_type == QuestionType.MCQ:
            matrix.mcq_count = num_questions
        elif question_type == QuestionType.TRUE_FALSE:
            matrix.true_false_count = num_questions
        elif question_type == QuestionType.SHORT_ANSWER:
            matrix.short_answer_count = num_questions
        else:
            matrix.mcq_count = max(1, num_questions // 3)
            matrix.true_false_count = max(1, num_questions // 3)
            matrix.short_answer_count = num_questions - matrix.mcq_count - matrix.true_false_count
        return self.generate_questions_matrix(text, matrix, user_notes)

    def _parse_questions_json(self, result: str) -> List[Dict[str, Any]]:
        result = result.strip()
        if '```json' in result:
            result = result.split('```json')[1].split('```')[0]
        elif '```' in result:
            result = result.split('```')[1].split('```')[0]
        result = result.strip()

        try:
            questions = json.loads(result)
            if isinstance(questions, list):
                return questions
            return []
        except json.JSONDecodeError as e:
            logger.error(f"Failed to parse questions JSON: {e}")
            return []

    def _fallback_questions(self, num_questions: int) -> List[Dict[str, Any]]:
        return [{
            'type': 'short_answer',
            'question': 'ما هي الفكرة الرئيسية في هذا النص؟',
            'answer': 'راجع النص للإجابة',
            'explanation': 'سؤال تلقائي - خدمة AI غير متاحة حالياً',
            'score': 1.0
        }]

    def ask_document(self, text: str, question: str, user_notes: str = "") -> str:
        """الإجابة على سؤال من سياق المستند."""
        chunks = self._chunker.chunk_text(text)
        config = _get_ai_config()
        input_limit = config.chunk_size if config else FALLBACK_CHUNK_SIZE
        context = chunks[0] if chunks else text[:input_limit]

        if len(chunks) > 1:
            context = self._find_relevant_chunks(chunks, question)

        notes_section = f"\nسياق إضافي: {user_notes}" if user_notes else ""

        prompt = f"""أنت مساعد أكاديمي يجيب على الأسئلة بناءً على محتوى المستندات.

قواعد:
1. أجب بناءً على المحتوى المقدم فقط
2. إذا لم تجد الإجابة، قل ذلك بوضوح
3. استخدم اللغة العربية الفصحى
4. كن واضحاً ومفصلاً
5. استخدم صيغة Markdown في الإجابة
{notes_section}

المحتوى:
{context}

السؤال: {question}

الإجابة:"""

        try:
            return self._generate_content(prompt, max_tokens=1000)
        except GeminiError as e:
            logger.error(f"Document Q&A failed: {e}")
            return "عذراً، حدث خطأ أثناء معالجة سؤالك. يرجى المحاولة مرة أخرى."

    def _find_relevant_chunks(self, chunks: List[str], question: str) -> str:
        question_words = set(question.lower().split())
        scored = []
        for chunk in chunks:
            chunk_words = set(chunk.lower().split())
            overlap = len(question_words & chunk_words)
            scored.append((overlap, chunk))
        scored.sort(key=lambda x: x[0], reverse=True)
        top_chunks = [c[1] for c in scored[:3]]
        return "\n\n---\n\n".join(top_chunks)

    def generate_and_save_summary(self, file_obj, user_notes: str = "") -> AIResponse:
        """توليد ملخص وحفظه كملف .md."""
        text = self.extract_text_from_file(file_obj)
        if not text:
            return AIResponse(success=False, error='لا يمكن استخراج النص من الملف')

        try:
            summary = self.generate_summary(text, user_notes=user_notes)
            md_path = self._storage.save_summary(
                file_id=file_obj.id,
                content=summary,
                metadata={
                    'source_file': file_obj.title,
                    'course': str(file_obj.course),
                    'date': datetime.now().strftime('%Y-%m-%d %H:%M'),
                    'model': self._model_name,
                }
            )
            return AIResponse(success=True, data=summary, md_file_path=md_path)
        except Exception as e:
            logger.error(f"generate_and_save_summary failed: {e}")
            return AIResponse(success=False, error=str(e))

    def generate_and_save_questions(
        self, file_obj, matrix: QuestionMatrixConfig, user_notes: str = ""
    ) -> AIResponse:
        """توليد أسئلة وحفظها كملف .md."""
        text = self.extract_text_from_file(file_obj)
        if not text:
            return AIResponse(success=False, error='لا يمكن استخراج النص من الملف')

        try:
            questions = self.generate_questions_matrix(text, matrix, user_notes)
            if not questions:
                return AIResponse(success=False, error='لم يتمكن AI من توليد أسئلة')

            md_path = self._storage.save_questions(
                file_id=file_obj.id,
                questions_data=questions,
                metadata={
                    'source_file': file_obj.title,
                    'course': str(file_obj.course),
                    'date': datetime.now().strftime('%Y-%m-%d %H:%M'),
                    'total_questions': str(len(questions)),
                    'total_score': str(matrix.total_score),
                    'model': self._model_name,
                }
            )
            return AIResponse(success=True, data=questions, md_file_path=md_path)
        except Exception as e:
            logger.error(f"generate_and_save_questions failed: {e}")
            return AIResponse(success=False, error=str(e))

    def test_connection(self) -> AIResponse:
        """اختبار الاتصال."""
        try:
            start_ms = int(time.time() * 1000)
            response = self._generate_content("قل: مرحباً، أنا جاهز!", max_tokens=50)
            latency = int(time.time() * 1000) - start_ms
            return AIResponse(success=True, data={'response': response, 'latency_ms': latency})
        except GeminiError as e:
            return AIResponse(success=False, error=str(e))


# ========== Celery Tasks (Optional) ==========

try:
    from celery import shared_task
    CELERY_AVAILABLE = True
except ImportError:
    CELERY_AVAILABLE = False
    def shared_task(*args, **kwargs):
        def decorator(func):
            return func
        return decorator


@shared_task(bind=True, max_retries=3, default_retry_delay=60)
def generate_summary_async(self, file_id: int, user_notes: str = "") -> Dict[str, Any]:
    from apps.courses.models import LectureFile
    from apps.ai_features.models import AISummary

    try:
        file_obj = LectureFile.objects.get(pk=file_id)
        service = GeminiService()
        result = service.generate_and_save_summary(file_obj, user_notes=user_notes)

        if result.success:
            AISummary.objects.update_or_create(
                file=file_obj,
                defaults={
                    'summary_text': result.data[:200] + '...' if len(result.data) > 200 else result.data,
                    'is_cached': True,
                    'model_used': service._model_name,
                }
            )
            return {'success': True, 'md_file_path': result.md_file_path}
        return {'success': False, 'error': result.error}

    except LectureFile.DoesNotExist:
        return {'success': False, 'error': 'الملف غير موجود'}
    except Exception as e:
        logger.error(f"Async summary generation failed: {e}")
        if CELERY_AVAILABLE and hasattr(self, 'retry'):
            raise self.retry(exc=e)
        return {'success': False, 'error': str(e)}


@shared_task(bind=True, max_retries=3, default_retry_delay=60)
def generate_questions_async(
    self, file_id: int, question_type: str = 'mixed',
    num_questions: int = 5, user_notes: str = ""
) -> Dict[str, Any]:
    from apps.courses.models import LectureFile
    from apps.ai_features.models import AIGeneratedQuestion

    try:
        file_obj = LectureFile.objects.get(pk=file_id)
        service = GeminiService()

        text = service.extract_text_from_file(file_obj)
        if not text:
            return {'success': False, 'error': 'لا يمكن استخراج النص'}

        q_type = QuestionType(question_type) if question_type in [e.value for e in QuestionType] else QuestionType.MIXED
        questions = service.generate_questions(text, q_type, num_questions, user_notes)

        saved_ids = []
        for q in questions:
            ai_q = AIGeneratedQuestion.objects.create(
                file=file_obj,
                question_type=q.get('type', 'short_answer'),
                question_text=q.get('question', ''),
                options=q.get('options'),
                correct_answer=q.get('answer', ''),
                explanation=q.get('explanation', ''),
            )
            saved_ids.append(ai_q.id)

        return {'success': True, 'question_ids': saved_ids, 'count': len(saved_ids)}

    except LectureFile.DoesNotExist:
        return {'success': False, 'error': 'الملف غير موجود'}
    except Exception as e:
        logger.error(f"Async question generation failed: {e}")
        if CELERY_AVAILABLE and hasattr(self, 'retry'):
            raise self.retry(exc=e)
        return {'success': False, 'error': str(e)}
